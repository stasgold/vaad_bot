from __future__ import annotations

import re
from datetime import datetime, timezone
from email import message_from_binary_file
from email.policy import default as email_policy
from pathlib import Path
from typing import Any, Dict, List, Tuple

import markdown
from mcp.server.fastmcp import FastMCP

mcp = FastMCP("doc2md")

ALL_MD_EXTENSIONS = [
    "extra",
    "tables",
    "fenced_code",
    "toc",
    "attr_list",
    "md_in_html",
    "sane_lists",
    "nl2br",
    "codehilite",
]


def _safe_text_read(file_path: Path) -> str:
    with file_path.open("rb") as f:
        raw = f.read()
    try:
        from charset_normalizer import from_bytes

        result = from_bytes(raw).best()
        if result is not None:
            return str(result)
    except Exception:
        pass
    return raw.decode("utf-8", errors="replace")


def _normalize_markdown(md_text: str) -> str:
    """
    Validate/normalize Markdown by running it through markdown[all].
    We return the original Markdown but ensure parsing succeeds.
    """
    markdown.markdown(md_text, extensions=ALL_MD_EXTENSIONS)
    return md_text.strip() + "\n"


def _md_heading(title: str, level: int = 1) -> str:
    hashes = "#" * max(1, min(6, level))
    return f"{hashes} {title}\n"


def _clean_text(text: str) -> str:
    text = text.replace("\u00a0", " ")
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def _extract_pdf(file_path: Path) -> Tuple[str, Dict[str, Any]]:
    import pdfplumber

    pages_text: List[str] = []
    tables_count = 0
    with pdfplumber.open(str(file_path)) as pdf:
        for i, page in enumerate(pdf.pages, start=1):
            page_text = page.extract_text() or ""
            tables = page.extract_tables() or []
            tables_count += len(tables)
            if page_text:
                pages_text.append(_md_heading(f"Page {i}", 2) + _clean_text(page_text))
    md = _md_heading(file_path.stem, 1) + "\n\n" + "\n\n".join(pages_text)
    return md.strip() + "\n", {"pages": len(pages_text), "tables": tables_count}


def _extract_docx(file_path: Path) -> Tuple[str, Dict[str, Any]]:
    from docx import Document

    doc = Document(str(file_path))
    parts: List[str] = [_md_heading(file_path.stem, 1)]
    tables_count = 0

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        style_name = (para.style.name or "").lower() if para.style else ""
        if style_name.startswith("heading"):
            level_match = re.findall(r"\d+", style_name)
            level = int(level_match[0]) if level_match else 2
            parts.append(_md_heading(text, level))
        else:
            parts.append(text)

    for table in doc.tables:
        tables_count += 1
        rows = []
        for row in table.rows:
            rows.append([cell.text.strip() for cell in row.cells])
        if rows:
            header = rows[0]
            sep = ["---"] * len(header)
            body = rows[1:] if len(rows) > 1 else []
            table_md = ["| " + " | ".join(header) + " |", "| " + " | ".join(sep) + " |"]
            for row in body:
                table_md.append("| " + " | ".join(row) + " |")
            parts.append("\n".join(table_md))

    md = "\n\n".join(p for p in parts if p)
    return md.strip() + "\n", {"tables": tables_count}


def _extract_pptx(file_path: Path) -> Tuple[str, Dict[str, Any]]:
    from pptx import Presentation

    prs = Presentation(str(file_path))
    parts: List[str] = [_md_heading(file_path.stem, 1)]
    for i, slide in enumerate(prs.slides, start=1):
        parts.append(_md_heading(f"Slide {i}", 2))
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                parts.append(shape.text.strip())
    md = "\n\n".join(p for p in parts if p)
    return md.strip() + "\n", {"slides": len(prs.slides)}


def _extract_excel(file_path: Path) -> Tuple[str, Dict[str, Any]]:
    import pandas as pd

    sheets = pd.read_excel(str(file_path), sheet_name=None)
    parts: List[str] = [_md_heading(file_path.stem, 1)]
    for name, df in sheets.items():
        parts.append(_md_heading(f"Sheet: {name}", 2))
        if df.empty:
            parts.append("(empty)")
        else:
            parts.append(df.to_markdown(index=False))
    md = "\n\n".join(parts)
    return md.strip() + "\n", {"sheets": len(sheets), "tables": len(sheets)}


def _extract_csv(file_path: Path) -> Tuple[str, Dict[str, Any]]:
    import pandas as pd

    df = pd.read_csv(str(file_path))
    parts = [_md_heading(file_path.stem, 1)]
    parts.append(df.to_markdown(index=False) if not df.empty else "(empty)")
    md = "\n\n".join(parts)
    return md.strip() + "\n", {"tables": 1}


def _extract_html(file_path: Path) -> Tuple[str, Dict[str, Any]]:
    from markdownify import markdownify as mdify

    html_text = _safe_text_read(file_path)
    md = mdify(html_text)
    return md.strip() + "\n", {}


def _extract_eml(file_path: Path) -> Tuple[str, Dict[str, Any]]:
    from markdownify import markdownify as mdify

    with file_path.open("rb") as f:
        msg = message_from_binary_file(f, policy=email_policy)

    parts: List[str] = []
    subject = msg.get("subject") or file_path.stem
    parts.append(_md_heading(subject, 1))

    for part in msg.walk():
        content_type = part.get_content_type()
        if content_type == "text/plain":
            text = part.get_content()
            parts.append(text.strip())
        elif content_type == "text/html":
            html = part.get_content()
            parts.append(mdify(html).strip())

    md = "\n\n".join(p for p in parts if p)
    return md.strip() + "\n", {}


def _extract_image(file_path: Path) -> Tuple[str, Dict[str, Any]]:
    from PIL import Image
    import pytesseract

    # Support Hebrew + English OCR
    # Use 'heb+eng' for Israeli documents with mixed Hebrew/English text
    try:
        text = pytesseract.image_to_string(
            Image.open(str(file_path)),
            lang='heb+eng'
        )
    except pytesseract.TesseractError:
        # Fallback to English only if Hebrew not available
        text = pytesseract.image_to_string(Image.open(str(file_path)))
    
    md = _md_heading(file_path.stem, 1) + "\n\n" + _clean_text(text)
    return md.strip() + "\n", {}


def _extract_text(file_path: Path) -> Tuple[str, Dict[str, Any]]:
    text = _safe_text_read(file_path)
    md = _md_heading(file_path.stem, 1) + "\n\n" + _clean_text(text)
    return md.strip() + "\n", {}


@mcp.tool()
def doc2md(file_path: str) -> Dict[str, Any]:
    """
    Convert a file into Markdown and return metadata.
    """
    path = Path(file_path).expanduser().resolve()
    if not path.exists() or not path.is_file():
        return {
            "error": "file_not_found",
            "message": f"File not found: {path}",
        }

    ext = path.suffix.lower()
    metadata: Dict[str, Any] = {
        "source_filename": path.name,
        "conversion_timestamp": datetime.now(timezone.utc).isoformat(),
        "document_type": ext.lstrip(".") or "unknown",
        "notes": "best-effort extraction",
    }

    try:
        if ext in {".pdf"}:
            md, extra = _extract_pdf(path)
        elif ext in {".docx"}:
            md, extra = _extract_docx(path)
        elif ext in {".pptx"}:
            md, extra = _extract_pptx(path)
        elif ext in {".xlsx", ".xls"}:
            md, extra = _extract_excel(path)
        elif ext in {".csv"}:
            md, extra = _extract_csv(path)
        elif ext in {".html", ".htm"}:
            md, extra = _extract_html(path)
        elif ext in {".eml"}:
            md, extra = _extract_eml(path)
        elif ext in {".png", ".jpg", ".jpeg", ".tiff", ".bmp", ".gif"}:
            md, extra = _extract_image(path)
        elif ext in {".md", ".txt"}:
            md, extra = _extract_text(path)
        else:
            md, extra = _extract_text(path)

        metadata.update(extra)
        md = _normalize_markdown(md)

        return {
            "markdown": md,
            "metadata": metadata,
        }
    except Exception as exc:
        return {
            "error": "conversion_failed",
            "message": str(exc),
            "metadata": metadata,
        }


if __name__ == "__main__":
    mcp.run()
